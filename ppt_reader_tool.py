from pptx import Presentation
from crewai.tools import tool


@tool("ppt_reader")
def ppt_reader(ppt_path: str):
    """
    Reads every slide and extracts raw content.
    """

    prs = Presentation(ppt_path)

    result = {
        "presentation_name": ppt_path,
        "total_slides": len(prs.slides),
        "slides": []
    }

    for slide_number, slide in enumerate(prs.slides, start=1):

        slide_info = {
            "slide_number": slide_number,
            "title": "",
            "texts": [],
            "tables": [],
            "images": [],
            "notes": ""
        }

        for shape in slide.shapes:

            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text == "":
                    continue

                if slide_info["title"] == "":
                    slide_info["title"] = text

                else:
                    slide_info["texts"].append(text)

        result["slides"].append(slide_info)

    return result