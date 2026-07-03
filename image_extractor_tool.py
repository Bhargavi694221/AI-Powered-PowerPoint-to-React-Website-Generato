from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import io

from crewai.tools import tool


@tool("Image Extractor")
def image_extractor(ppt_file: str) -> str:
    """
    Extract all embedded images from a PowerPoint presentation
    and save them into generated_website/assets.
    """

    ppt_path = Path(ppt_file)

    if not ppt_path.exists():
        return f"❌ PowerPoint not found: {ppt_file}"

    output_folder = Path("generated_website/assets")
    output_folder.mkdir(parents=True, exist_ok=True)

    prs = Presentation(ppt_path)

    image_count = 0

    for slide_number, slide in enumerate(prs.slides, start=1):

        for shape in slide.shapes:

            # Only process picture shapes
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            try:
                image = shape.image

                image_bytes = io.BytesIO(image.blob)

                img = Image.open(image_bytes)

                extension = image.ext.lower()

                filename = output_folder / (
                    f"slide_{slide_number}_{image_count + 1}.{extension}"
                )

                img.save(filename)

                image_count += 1

            except Exception as e:
                print(f"Skipping image on slide {slide_number}: {e}")

    return (
        f"✅ Successfully extracted {image_count} images.\n"
        f"Saved in: {output_folder}"
    )