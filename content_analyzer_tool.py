import json
from pathlib import Path
from pptx import Presentation
from crewai.tools import tool


@tool("Content Analyzer")
def content_analyzer(ppt_file: str) -> str:
    """
    Reads a PowerPoint presentation and converts it into
    structured JSON for AI website generation.
    """

    ppt_path = Path(ppt_file)

    if not ppt_path.exists():
        return f"❌ PowerPoint not found: {ppt_file}"

    prs = Presentation(ppt_path)

    website_content = {
        "project_name": "",
        "slides": []
    }

    for slide_number, slide in enumerate(prs.slides, start=1):

        slide_data = {
            "slide_number": slide_number,
            "title": "",
            "content": [],
            "image_count": 0
        }

        for shape in slide.shapes:

            # Extract text
            if hasattr(shape, "text"):

                text = shape.text.strip()

                if text:

                    if slide_data["title"] == "":
                        slide_data["title"] = text

                        if website_content["project_name"] == "":
                            website_content["project_name"] = text

                    else:
                        slide_data["content"].append(text)

            # Count images
            if hasattr(shape, "image"):
                slide_data["image_count"] += 1

        if slide_data["title"] == "":
            slide_data["title"] = f"Slide {slide_number}"

        website_content["slides"].append(slide_data)

    output_folder = Path("generated_website/data")
    output_folder.mkdir(parents=True, exist_ok=True)

    json_file = output_folder / "website_content.json"

    with open(json_file, "w", encoding="utf-8") as f:
        json.dump(
            website_content,
            f,
            indent=4,
            ensure_ascii=False
        )

    return (
        f"✅ Content analysis completed.\n"
        f"Slides analyzed : {len(website_content['slides'])}\n"
        f"Output : {json_file}"
    )